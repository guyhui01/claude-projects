"""
Génération Story Map PowerPoint — Company-Test e-Commerce B2B
Méthode : Jeff Patton Story Mapping
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── PALETTE ───────────────────────────────────────────────────────────────────
DARK_NAVY    = RGBColor(0x1A, 0x1A, 0x2E)
NAVY         = RGBColor(0x16, 0x21, 0x3E)
MUST         = RGBColor(0xFF, 0xC1, 0x07)   # jaune Must Have
MUST_TEXT    = RGBColor(0x85, 0x64, 0x04)
SHOULD       = RGBColor(0x17, 0xA2, 0xB8)   # bleu Should
SHOULD_BG    = RGBColor(0xD1, 0xEC, 0xF1)
COULD        = RGBColor(0x28, 0xA7, 0x45)   # vert Could
COULD_BG     = RGBColor(0xD4, 0xED, 0xDA)
PHASE2       = RGBColor(0xDC, 0x35, 0x45)   # rouge Phase 2
PHASE2_BG    = RGBColor(0xF8, 0xD7, 0xDA)
DONE_BG      = RGBColor(0xC3, 0xE6, 0xCB)   # vert clair Livré
DONE_BORDER  = RGBColor(0x15, 0x57, 0x24)
SPIKE_BG     = RGBColor(0xE2, 0xD9, 0xF3)
SPIKE_BORDER = RGBColor(0x6F, 0x42, 0xC1)
BACKBONE_BG  = RGBColor(0x1A, 0x1A, 0x2E)
TASK_BG      = RGBColor(0xE8, 0xEC, 0xF0)
SWIM_BG      = RGBColor(0x49, 0x50, 0x57)
RELEASE_S1   = RGBColor(0xFF, 0xF3, 0xCD)
RELEASE_S2   = RGBColor(0xD1, 0xEC, 0xF1)
RELEASE_S3   = RGBColor(0xCC, 0xE5, 0xFF)
RELEASE_S4   = RGBColor(0xD4, 0xED, 0xDA)
RELEASE_IP   = RGBColor(0xE2, 0xD9, 0xF3)
RELEASE_PH2  = RGBColor(0xF8, 0xD7, 0xDA)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF4, 0xF6, 0xF9)

# ── HELPERS ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=8, bold=False,
                font_color=RGBColor(0x1A,0x1A,0x2E), align=PP_ALIGN.CENTER,
                fill_color=None, border_color=None, border_width=Pt(0.5),
                wrap=True, v_anchor=None):
    if fill_color:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = border_width
        else:
            shape.line.fill.background()
        tf = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame

    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor

    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    return shape

def card(slide, x, y, w, h, us_id, label, bg, border, done=False):
    """Crée une carte User Story."""
    shape = add_rect(slide, x, y, w, h, bg, border, Pt(1.5))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    # ID
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = us_id + (" ✓" if done else "")
    r1.font.size = Pt(6.5)
    r1.font.bold = True
    r1.font.color.rgb = border
    # Label
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(7.5)
    r2.font.bold = False
    r2.font.color.rgb = DARK_NAVY
    return shape

def release_band(slide, x, y, w, h, label, bg, border):
    shape = add_rect(slide, x, y, w, h, bg, border, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = border
    return shape

# ══════════════════════════════════════════════════════════════════════════════
# PRÉSENTATION
# ══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(16.54)   # A3 paysage largeur
prs.slide_height = Inches(11.69)   # A3 paysage hauteur

blank = prs.slide_layouts[6]       # Layout vide

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank)

# Fond
bg = add_rect(slide1, 0, 0, 16.54, 11.69, DARK_NAVY)

# Bande couleur gauche
add_rect(slide1, 0, 0, 0.35, 11.69, RGBColor(0xF0, 0xA5, 0x00))

# Titre principal
add_textbox(slide1, 1, 3.5, 14, 1.5,
            "Story Map\nCompany-Test e-Commerce B2B",
            font_size=36, bold=True, font_color=WHITE, align=PP_ALIGN.LEFT)

# Sous-titre
add_textbox(slide1, 1, 5.3, 10, 0.8,
            "Drupal 10 + Commerce 2.x   ·   PI-2026-Q3   ·   23 User Stories   ·   70 SP",
            font_size=14, bold=False, font_color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT)

# Méta
add_textbox(slide1, 1, 7, 8, 1.2,
            "Atelier Story Mapping — 2026-04-15\nParticipants : PO · BA · UX/UI · Dev Lead\nMéthode : Jeff Patton Story Mapping",
            font_size=11, bold=False, font_color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.LEFT)

# Logo/badge
add_textbox(slide1, 12.5, 9.5, 3, 1,
            "Company-Test\ne-Commerce B2B",
            font_size=10, bold=True, font_color=RGBColor(0xF0,0xA5,0x00), align=PP_ALIGN.RIGHT)

# Légende
legends = [
    ("Must Have",    MUST,      RGBColor(0x85,0x64,0x04)),
    ("Should Have",  SHOULD_BG, SHOULD),
    ("Could Have",   COULD_BG,  COULD),
    ("Phase 2",      PHASE2_BG, PHASE2),
    ("Livré ✓",      DONE_BG,   DONE_BORDER),
    ("Spike",        SPIKE_BG,  SPIKE_BORDER),
]
for i, (lbl, bg_c, border_c) in enumerate(legends):
    xi = 1 + i * 2.4
    add_rect(slide1, xi, 9.2, 0.25, 0.25, bg_c, border_c, Pt(1.5))
    add_textbox(slide1, xi + 0.3, 9.15, 2, 0.35, lbl,
                font_size=9, bold=False, font_color=WHITE, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — STORY MAP : SWIMLANE ACHETEUR B2B
# ══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank)
add_rect(slide2, 0, 0, 16.54, 11.69, LIGHT_GRAY)

# Titre slide
add_textbox(slide2, 0.2, 0.1, 16, 0.45,
            "Story Map — Swimlane Acheteur B2B   |   Company-Test e-Commerce B2B   |   PI-2026-Q3",
            font_size=10, bold=True, font_color=WHITE, align=PP_ALIGN.LEFT,
            fill_color=DARK_NAVY)

# ── SWIMLANE LABEL ──
add_textbox(slide2, 0.2, 0.65, 15.9, 0.35,
            "👤  SWIMLANE — ACHETEUR B2B",
            font_size=10, bold=True, font_color=WHITE,
            fill_color=SWIM_BG, align=PP_ALIGN.LEFT)

# ── Colonnes définition ──
# Activités : 6 colonnes principales
acts = [
    (0.20, 1.85, "🔍 Découvrir\nle catalogue"),
    (3.00, 1.85, "📝 S'inscrire\ncompte B2B"),
    (4.85, 2.70, "🔑 Se connecter"),
    (7.55, 0.00, ""),   # colonne vide (fusionnée)
    (7.55, 3.60, "🛒 Commander"),
    (11.15, 1.85, "💳 Payer"),
    (13.00, 1.85, "📦 Suivre"),
]

# ── BACKBONE ──
backbones = [
    (0.20, 1.80, "🔍 Découvrir\nle catalogue"),
    (2.05, 1.80, "📝 S'inscrire\ncompte B2B"),
    (3.90, 2.65, "🔑 Se connecter"),
    (6.60, 3.55, "🛒 Commander"),
    (10.20, 1.85, "💳 Payer"),
    (12.10, 3.40, "📦 Suivre"),
]
for (bx, bw, bl) in backbones:
    shape = add_rect(slide2, bx, 1.1, bw, 0.7, BACKBONE_BG, NAVY, Pt(0.5))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = bl
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = WHITE

# ── TÂCHES ──
tasks = [
    (0.20, 1.80, "Naviguer\npar catégorie"),
    (2.05, 1.80, "Formulaire\nd'inscription"),
    (3.90, 1.30, "Connexion\nemail + mdp"),
    (5.25, 1.30, "Gestion\nstatuts compte"),
    (6.60, 1.75, "Panier &\nN° BC"),
    (8.40, 1.75, "Validation\ncommande"),
    (10.20, 1.85, "Paiement\nCB 3DS"),
    (12.10, 1.70, "Statut &\nsuivi"),
    (13.85, 1.65, "Historique\ncommandes"),
]
for (tx, tw, tl) in tasks:
    shape = add_rect(slide2, tx, 1.90, tw, 0.6, TASK_BG,
                     RGBColor(0x6C,0x75,0x7D), Pt(0.5))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tl
    r.font.size = Pt(7.5)
    r.font.bold = True
    r.font.color.rgb = DARK_NAVY

# ── RELEASE BANDS + CARTES ──
CW = 1.72   # card width
CH = 0.75   # card height
GAP = 0.08  # gap between cards

# Sprint 1
release_band(slide2, 0.20, 2.62, 15.5, 0.28,
             "  🏃 Sprint 1 — Walking Skeleton : compte B2B activé + catalogue navigable",
             RELEASE_S1, RGBColor(0xF0,0xA5,0x00))

s1_cards = [
    (0.20,  "US-001", "Navigation\npar catégorie",   DONE_BG, DONE_BORDER, True),
    (1.97,  "US-002", "Fiche produit\n+ tarif B2B",  DONE_BG, DONE_BORDER, True),
    (3.74,  "US-005", "Formulaire\n+ SIRET regex",   DONE_BG, DONE_BORDER, True),
    (5.51,  "US-006", "Connexion\nréussie",          DONE_BG, DONE_BORDER, True),
    (7.28,  "US-006", "Statuts: att.\nrefusé, bloqué",DONE_BG,DONE_BORDER, True),
    (10.20, "SPIKE",  "Éval. PSP\nStripe ✓",         SPIKE_BG, SPIKE_BORDER, False),
]
for (cx, uid, lbl, bg_c, bd_c, dn) in s1_cards:
    card(slide2, cx, 2.95, CW, CH, uid, lbl, bg_c, bd_c, dn)

# Sprint 2
release_band(slide2, 0.20, 3.82, 15.5, 0.28,
             "  🏃 Sprint 2 — Commande bout en bout",
             RELEASE_S2, SHOULD)

s2_cards = [
    (6.60,  "US-009", "Ajout au\npanier",            MUST, MUST_TEXT, False),
    (8.37,  "US-010", "N° bon de\ncommande",         MUST, MUST_TEXT, False),
    (10.14, "US-011", "Email\nconfirmation",          MUST, MUST_TEXT, False),
]
for (cx, uid, lbl, bg_c, bd_c, dn) in s2_cards:
    card(slide2, cx, 4.15, CW, CH, uid, lbl, bg_c, bd_c, dn)

# Sprint 3
release_band(slide2, 0.20, 5.02, 15.5, 0.28,
             "  🏃 Sprint 3 — Paiement & Facturation",
             RELEASE_S3, RGBColor(0x00,0x7B,0xFF))

s3_cards = [
    (10.20, "US-016", "Paiement CB\n3D Secure",      MUST, MUST_TEXT, False),
    (11.97, "US-012", "Historique\ncommandes",        SHOULD_BG, SHOULD, False),
]
for (cx, uid, lbl, bg_c, bd_c, dn) in s3_cards:
    card(slide2, cx, 5.35, CW, CH, uid, lbl, bg_c, bd_c, dn)

# Sprint 4
release_band(slide2, 0.20, 6.22, 15.5, 0.28,
             "  🏃 Sprint 4 — Enrichissement UX & Livraisons",
             RELEASE_S4, COULD)

s4_cards = [
    (0.20,  "US-003", "Filtres\nproduits",           SHOULD_BG, SHOULD, False),
    (13.78, "US-019", "Frais livraison\nautomatique", SHOULD_BG, SHOULD, False),
    (12.05, "US-021", "N° suivi\ntransporteur",       SHOULD_BG, SHOULD, False),
]
for (cx, uid, lbl, bg_c, bd_c, dn) in s4_cards:
    card(slide2, cx, 6.55, CW, CH, uid, lbl, bg_c, bd_c, dn)

# IP Sprint
release_band(slide2, 0.20, 7.42, 15.5, 0.28,
             "  🔬 IP Sprint — Could Have",
             RELEASE_IP, SPIKE_BORDER)

ip_cards = [
    (0.20,  "US-004", "Fiche tech.\nPDF",            COULD_BG, COULD, False),
    (3.74,  "US-023", "Email bienvenue\n+ guide",     COULD_BG, COULD, False),
    (13.78, "US-020", "Adresses\nlivraison mult.",    COULD_BG, COULD, False),
]
for (cx, uid, lbl, bg_c, bd_c, dn) in ip_cards:
    card(slide2, cx, 7.75, CW, CH, uid, lbl, bg_c, bd_c, dn)

# Phase 2
release_band(slide2, 0.20, 8.62, 15.5, 0.28,
             "  🚀 Phase 2 — Hors scope V1",
             RELEASE_PH2, PHASE2)

ph2_cards = [
    (3.74,  "—",      "SIRET via\nAPI INSEE",        PHASE2_BG, PHASE2, False),
    (10.20, "US-017", "Paiement\ndifféré 30j",        PHASE2_BG, PHASE2, False),
]
for (cx, uid, lbl, bg_c, bd_c, dn) in ph2_cards:
    card(slide2, cx, 8.95, CW, CH, uid, lbl, bg_c, bd_c, dn)

# Footer
add_textbox(slide2, 0.20, 10.85, 15.9, 0.35,
            "Company-Test e-Commerce B2B  ·  PI-2026-Q3  ·  Vélocité Sprint 1 : 16 SP  ·  Atelier Story Mapping 2026-04-15",
            font_size=7, bold=False, font_color=WHITE,
            fill_color=DARK_NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — STORY MAP : SWIMLANE ADMIN Company-Test
# ══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank)
add_rect(slide3, 0, 0, 16.54, 11.69, LIGHT_GRAY)

add_textbox(slide3, 0.2, 0.1, 16, 0.45,
            "Story Map — Swimlane Admin Company-Test   |   Company-Test e-Commerce B2B   |   PI-2026-Q3",
            font_size=10, bold=True, font_color=WHITE, align=PP_ALIGN.LEFT,
            fill_color=DARK_NAVY)

add_textbox(slide3, 0.2, 0.65, 15.9, 0.35,
            "🔧  SWIMLANE — ADMIN Company-Test",
            font_size=10, bold=True, font_color=WHITE,
            fill_color=SWIM_BG, align=PP_ALIGN.LEFT)

# Backbone Admin
backbones_admin = [
    (0.20, 7.25, "👥 Gérer les comptes B2B"),
    (7.55, 7.25, "📋 Gérer les commandes"),
]
for (bx, bw, bl) in backbones_admin:
    shape = add_rect(slide3, bx, 1.1, bw, 0.7, BACKBONE_BG, NAVY, Pt(0.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = bl
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = WHITE

# Tâches Admin
tasks_admin = [
    (0.20,  1.80, "Notif\nnouvelle demande"),
    (2.05,  1.80, "Liste demandes\nen attente"),
    (3.90,  1.80, "Valider\ncompte"),
    (5.75,  1.80, "Rejeter\navec motif"),
    (7.60,  1.80, "Dashboard\ncommandes"),
    (9.45,  1.80, "Statut +\nnotification"),
    (11.30, 1.80, "Export\nCSV"),
    (13.15, 1.80, "Facture\nPDF"),
]
for (tx, tw, tl) in tasks_admin:
    shape = add_rect(slide3, tx, 1.90, tw, 0.6, TASK_BG,
                     RGBColor(0x6C,0x75,0x7D), Pt(0.5))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tl
    r.font.size = Pt(7.5)
    r.font.bold = True
    r.font.color.rgb = DARK_NAVY

# Sprint 1→2 (US-007)
release_band(slide3, 0.20, 2.62, 15.5, 0.28,
             "  🏃 Sprint 1 → Sprint 2 (US-007 reportée) — Validation des comptes B2B",
             RELEASE_S1, RGBColor(0xF0,0xA5,0x00))

s1_admin = [
    (0.20, "US-007", "Notif admin\nnouvelle demande", MUST, MUST_TEXT),
    (2.05, "US-007", "Badge compteur\nen attente",     MUST, MUST_TEXT),
    (3.90, "US-007", "Valider compte\n→ email activ.", MUST, MUST_TEXT),
    (5.75, "US-007", "Rejeter + motif\nobligatoire",   MUST, MUST_TEXT),
]
for (cx, uid, lbl, bg_c, bd_c) in s1_admin:
    card(slide3, cx, 2.95, CW, CH, uid, lbl, bg_c, bd_c, False)

# Sprint 2 Admin
release_band(slide3, 0.20, 3.82, 15.5, 0.28,
             "  🏃 Sprint 2 — Gestion des commandes",
             RELEASE_S2, SHOULD)

s2_admin = [
    (7.60,  "US-013", "Dashboard\ncommandes B2B", MUST, MUST_TEXT),
    (9.45,  "US-014", "Statut +\nnotif client",   MUST, MUST_TEXT),
]
for (cx, uid, lbl, bg_c, bd_c) in s2_admin:
    card(slide3, cx, 4.15, CW, CH, uid, lbl, bg_c, bd_c, False)

# Sprint 3 Admin
release_band(slide3, 0.20, 5.02, 15.5, 0.28,
             "  🏃 Sprint 3 — Facturation",
             RELEASE_S3, RGBColor(0x00,0x7B,0xFF))

s3_admin = [
    (11.30, "US-015", "Export\ncommandes CSV", SHOULD_BG, SHOULD),
    (13.15, "US-018", "Génération\nfacture PDF",  SHOULD_BG, SHOULD),
]
for (cx, uid, lbl, bg_c, bd_c) in s3_admin:
    card(slide3, cx, 5.35, CW, CH, uid, lbl, bg_c, bd_c, False)

# IP Sprint Admin
release_band(slide3, 0.20, 6.22, 15.5, 0.28,
             "  🔬 IP Sprint — Could Have",
             RELEASE_IP, SPIKE_BORDER)

ip_admin = [
    (0.20, "US-022", "Notif email\ndemande compte", COULD_BG, COULD),
    (13.15,"US-008", "Tarifs négociés\npar client",  COULD_BG, COULD),
]
for (cx, uid, lbl, bg_c, bd_c) in ip_admin:
    card(slide3, cx, 6.55, CW, CH, uid, lbl, bg_c, bd_c, False)

# Footer
add_textbox(slide3, 0.20, 10.85, 15.9, 0.35,
            "Company-Test e-Commerce B2B  ·  PI-2026-Q3  ·  Atelier Story Mapping 2026-04-15",
            font_size=7, bold=False, font_color=WHITE,
            fill_color=DARK_NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TRAÇABILITÉ US → SPRINT → MoSCoW
# ══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank)
add_rect(slide4, 0, 0, 16.54, 11.69, LIGHT_GRAY)

add_textbox(slide4, 0.2, 0.1, 16, 0.45,
            "Traçabilité Story Map → Backlog   |   Company-Test e-Commerce B2B",
            font_size=10, bold=True, font_color=WHITE,
            fill_color=DARK_NAVY, align=PP_ALIGN.LEFT)

headers = ["US ID", "Titre", "Activité", "Sprint", "MoSCoW", "SP", "Statut"]
col_widths = [0.8, 3.5, 2.5, 0.9, 1.3, 0.6, 1.5]
col_colors = [DARK_NAVY]*7

rows_data = [
    ("US-001", "Navigation catalogue par catégorie",   "Découvrir",    "S1", "Must",    "3",  "✅ Livré"),
    ("US-002", "Fiche produit complète + tarif B2B",   "Découvrir",    "S1", "Must",    "5",  "✅ Livré"),
    ("US-003", "Filtres produits",                     "Découvrir",    "S4", "Should",  "3",  "À faire"),
    ("US-004", "Fiche technique PDF",                  "Découvrir",    "IP", "Could",   "2",  "À faire"),
    ("US-005", "Création compte professionnel B2B",    "S'inscrire",   "S1", "Must",    "5",  "✅ Livré"),
    ("US-006", "Connexion acheteur B2B",               "Se connecter", "S1", "Must",    "2",  "✅ Livré"),
    ("US-007", "Validation manuelle compte (admin)",   "Admin comptes","S2", "Must",    "3",  "⚠️ Reporté S2"),
    ("US-008", "Tarifs négociés par client",           "Admin comptes","IP", "Could",   "5",  "À faire"),
    ("US-009", "Ajout produits au panier",             "Commander",    "S2", "Must",    "3",  "🔄 En cours"),
    ("US-010", "Saisie N° bon de commande",            "Commander",    "S2", "Must",    "2",  "🔄 En cours"),
    ("US-011", "Email confirmation commande",          "Commander",    "S2", "Must",    "2",  "🔄 En cours"),
    ("US-012", "Historique commandes acheteur",        "Suivre",       "S3", "Should",  "3",  "À faire"),
    ("US-013", "Dashboard back-office commandes",      "Admin cmd",    "S2", "Must",    "3*", "🔄 En cours"),
    ("US-014", "Mise à jour statut + notif client",    "Admin cmd",    "S2", "Must",    "3",  "🔄 En cours"),
    ("US-015", "Export commandes CSV",                 "Admin cmd",    "S3", "Should",  "2",  "À faire"),
    ("US-016", "Paiement CB 3D Secure (Stripe)",       "Payer",        "S3", "Must",    "8",  "À faire"),
    ("US-017", "Paiement différé 30j",                 "Payer",        "Ph2","Phase 2", "5",  "Hors scope"),
    ("US-018", "Génération facture PDF",               "Admin cmd",    "S3", "Should",  "3",  "À faire"),
    ("US-019", "Calcul frais livraison automatique",   "Suivre",       "S4", "Should",  "5",  "À faire"),
    ("US-020", "Adresses livraison multiples",         "Suivre",       "IP", "Could",   "3",  "À faire"),
    ("US-021", "Envoi numéro suivi transporteur",      "Suivre",       "S4", "Should",  "3",  "À faire"),
    ("US-022", "Notification email demande compte",    "Admin comptes","IP", "Could",   "2",  "À faire"),
    ("US-023", "Email bienvenue + guide démarrage",    "S'inscrire",   "IP", "Could",   "2",  "À faire"),
]

STATUS_COLORS = {
    "✅ Livré":       DONE_BG,
    "🔄 En cours":    RELEASE_S2,
    "⚠️ Reporté S2": RELEASE_S1,
    "Hors scope":     PHASE2_BG,
}
MOSCOW_COLORS = {
    "Must":    MUST,
    "Should":  SHOULD_BG,
    "Could":   COULD_BG,
    "Phase 2": PHASE2_BG,
}

# Header row
x_start = 0.20
y_header = 0.70
row_h = 0.36
x = x_start
for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
    add_textbox(slide4, x, y_header, cw, row_h-0.04, hdr,
                font_size=8, bold=True, font_color=WHITE,
                fill_color=DARK_NAVY, align=PP_ALIGN.CENTER)
    x += cw + 0.06

# Data rows
for ri, row in enumerate(rows_data):
    y = y_header + row_h + ri * row_h
    x = x_start
    for ci, (val, cw) in enumerate(zip(row, col_widths)):
        if ci == 4:   # MoSCoW
            bg_c = MOSCOW_COLORS.get(val, LIGHT_GRAY)
            fc   = DARK_NAVY
        elif ci == 6:  # Statut
            bg_c = STATUS_COLORS.get(val, LIGHT_GRAY)
            fc   = DARK_NAVY
        elif ri % 2 == 0:
            bg_c = WHITE
            fc   = DARK_NAVY
        else:
            bg_c = RGBColor(0xF8,0xF9,0xFA)
            fc   = DARK_NAVY
        add_textbox(slide4, x, y, cw, row_h-0.04, val,
                    font_size=7.5, bold=(ci==0), font_color=fc,
                    fill_color=bg_c, align=PP_ALIGN.CENTER,
                    border_color=RGBColor(0xDE,0xE2,0xE6))
        x += cw + 0.06

add_textbox(slide4, 0.20, 10.85, 15.9, 0.35,
            "* US-013 réévaluée 5 → 3 SP après Sprint 1  ·  Total V1 : 70 SP · 23 US  ·  PI-2026-Q3",
            font_size=7, bold=False, font_color=WHITE,
            fill_color=DARK_NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DÉCISIONS ATELIER
# ══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank)
add_rect(slide5, 0, 0, 16.54, 11.69, LIGHT_GRAY)

add_textbox(slide5, 0.2, 0.1, 16, 0.45,
            "Décisions clés — Atelier Story Mapping 2026-04-15   |   Company-Test e-Commerce B2B",
            font_size=10, bold=True, font_color=WHITE,
            fill_color=DARK_NAVY, align=PP_ALIGN.LEFT)

decisions = [
    ("D1", "Tarif B2B visible dès la fiche produit pour tout acheteur connecté",
           "Critère go/no-go recette finale", "PO + Direction Commerciale", MUST),
    ("D2", "Story map en 2 swimlanes : Acheteur B2B / Admin Company-Test",
           "Clarté du parcours dual", "PO + BA", SHOULD_BG),
    ("D3", "Validation SIRET = regex 14 chiffres V1 — API INSEE Phase 2",
           "US-005 — Réduction complexité Sprint 1", "PO + BA + DEV", MUST),
    ("D4", "Fiche technique PDF = Could Have (bouton conditionnel à l'existence du fichier)",
           "US-004 — Catalogue pas encore numérisé à 100%", "PO + UX", COULD_BG),
    ("D5", "N° bon de commande = Must Have, champ obligatoire au checkout",
           "US-010 — Obligation comptable clients Company-Test", "PO + BA", MUST),
    ("D6", "SPIKE PSP Sprint 1 (1j timebox) → Stripe retenu le 2026-05-20",
           "US-016 — Confirmé après évaluation Stripe vs PayPlug", "PO + DEV", SPIKE_BG),
    ("D7", "Paiement différé 30j = Phase 2 — hors scope V1",
           "US-017 — Complexité juridique + intégration ERP", "PO + DEV + BA", PHASE2_BG),
    ("D8", "Walking skeleton Sprint 1 = compte B2B activé + catalogue navigable",
           "Milestone M1 — Démo interne Direction Commerciale", "Équipe", DONE_BG),
    ("D9", "Wireframes catalogue + fiche produit = critère DoR Sprint 1",
           "Livraison UX avant Sprint Planning", "PO + UX", RELEASE_S1),
    ("D10","BA rédige les 4 templates email avant Sprint Planning Sprint 1",
           "Emails : confirmation, notif admin, activation, refus", "BA", SHOULD_BG),
]

for i, (did, decision, impact, owner, color) in enumerate(decisions):
    col = i % 2
    row = i // 2
    x = 0.20 + col * 8.0
    y = 0.75 + row * 2.0

    shape = add_rect(slide5, x, y, 7.8, 1.85, color,
                     RGBColor(0xDE,0xE2,0xE6), Pt(1))

    # Badge ID
    add_textbox(slide5, x+0.1, y+0.1, 0.55, 0.45, did,
                font_size=9, bold=True, font_color=WHITE,
                fill_color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # Décision
    add_textbox(slide5, x+0.75, y+0.08, 6.9, 0.55, decision,
                font_size=9, bold=True, font_color=DARK_NAVY, align=PP_ALIGN.LEFT)

    # Impact
    add_textbox(slide5, x+0.75, y+0.65, 6.9, 0.45,
                "Impact : " + impact,
                font_size=8, bold=False, font_color=RGBColor(0x49,0x50,0x57),
                align=PP_ALIGN.LEFT)

    # Owner
    add_textbox(slide5, x+0.75, y+1.1, 6.9, 0.45,
                "Acté par : " + owner,
                font_size=8, bold=False, font_color=RGBColor(0x6C,0x75,0x7D),
                align=PP_ALIGN.LEFT)

add_textbox(slide5, 0.20, 10.85, 15.9, 0.35,
            "CR complet : 01_po_scrum/story_mapping/cr_story_mapping.md  ·  Company-Test e-Commerce B2B  ·  PI-2026-Q3",
            font_size=7, bold=False, font_color=WHITE,
            fill_color=DARK_NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
output = r"C:\Users\Guy HUIBONHOA\ClaudeCode\projects\Project_test__ecommerce_b2b\01_po_scrum\story_mapping\story_map_Company-Test_b2b.pptx"
prs.save(output)
print(f"✅ Fichier généré : {output}")
